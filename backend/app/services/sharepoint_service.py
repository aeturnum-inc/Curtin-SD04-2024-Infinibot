import io
import requests
from typing import List, Dict, Any, Optional
import datetime

from azure.identity import ClientSecretCredential
import pdfplumber
import docx2txt
import openpyxl
import csv
from pptx import Presentation 
from pymongo import MongoClient
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from app.services.embedding_service import get_embedding_model
from app.services.document_permission import get_document_permissions
from langchain_mongodb import MongoDBAtlasVectorSearch

from app.core.config import settings


class SharePointService:
    """
    Service for interacting with SharePoint via Microsoft Graph API.
    """
    
    def __init__(self):
        """
        Initialize the SharePoint service with Azure credentials.
        """
        self.credential = ClientSecretCredential(
            tenant_id=settings.TENANT_ID,
            client_id=settings.CLIENT_ID,
            client_secret=settings.CLIENT_SECRET
        )
        self.access_token = None
        
         # Initialize MongoDB client
        self.mongo_client = MongoClient(settings.MONGODB_ATLAS_URI)
        self.db = self.mongo_client[settings.DB_NAME]
        self.collection = self.db[settings.COLLECTION_NAME]


    def get_access_token(self) -> str:
        """
        Get a valid access token for Microsoft Graph API.
        
        Returns:
            str: The access token
        """
        if not self.access_token:
            token_response = self.credential.get_token("https://graph.microsoft.com/.default")
            self.access_token = token_response.token
        return self.access_token

    def list_drives(self) -> List[Dict[str, Any]]:
        """
        List all drives in the SharePoint site.

        Returns:
            List[Dict[str, Any]]: List of drives
        """
        endpoint = f"https://graph.microsoft.com/v1.0/sites/{settings.SITE_ID}/drives"
        headers = {"Authorization": f"Bearer {self.get_access_token()}"}

        response = requests.get(endpoint, headers=headers)
        response.raise_for_status()
        return response.json().get("value", [])

    def list_documents(self, user_email: str = None) -> List[Dict[str, Any]]:
        """
        List all documents in the SharePoint site that the user has access to.
        
        Args:
            user_email: The email of the user to check permissions for
            
        Returns:
            List[Dict[str, Any]]: List of document metadata
        """
        token = self.get_access_token()
        drives_endpoint = f"https://graph.microsoft.com/v1.0/sites/{settings.SITE_ID}/drives"
        headers = {"Authorization": f"Bearer {token}"}
        
        # Get drives
        drives_response = requests.get(drives_endpoint, headers=headers)
        drives_response.raise_for_status()
        drives = drives_response.json().get("value", [])
        
        if not drives:
            raise Exception("No drives found in SharePoint response")
        
        drive_id = drives[0]["id"]
        
        # Get documents
        documents_endpoint = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
        documents_response = requests.get(documents_endpoint, headers=headers)
        documents_response.raise_for_status()
        documents = documents_response.json().get("value", [])
        
        if not documents:
            raise Exception("No documents found in SharePoint response")
        
        # If user_email is provided, filter documents based on permissions
        if user_email and not settings.DEV_MODE:
            accessible_documents = []
            for doc in documents:
                if self.check_user_permission(doc["id"], user_email, drive_id):
                    accessible_documents.append({
                        "id": doc["id"],
                        "name": doc.get("name", "Untitled"),
                        "webUrl": doc["webUrl"],
                        "lastModified": doc["lastModifiedDateTime"],
                    })
            return accessible_documents
        else:
            # In dev mode or without user_email, return all documents
            return [
                {
                    "id": item["id"],
                    "name": item.get("name", "Untitled"),
                    "webUrl": item["webUrl"],
                    "lastModified": item["lastModifiedDateTime"],
                }
                for item in documents
            ]

    def check_user_permission(self, document_id: str, user_email: str, drive_id: str) -> bool:
        """
        Enhanced check if a user has permission to access a document.
        Supports users, groups, organization-wide, and public access.
        
        Args:
            document_id: The document ID
            user_email: The user's email
            drive_id: The drive ID
            
        Returns:
            bool: True if the user has access, False otherwise
        """
        try:
            # Get comprehensive permissions for the document
            permission_data = get_document_permissions(self, document_id, drive_id)
            
            # Check access level first
            access_level = permission_data.get("access_level", "private")
            
            # Public access - anyone can access
            if access_level == "public":
                print(f"Document has public access - allowing access to {user_email}")
                return True
            
            # Organization access - check if user is in the organization
            if access_level == "organization":
                from app.services.document_permission import is_user_in_organization
                if is_user_in_organization(user_email, self):
                    print(f"Document has organization access and user is in org - allowing access to {user_email}")
                    return True
                else:
                    print(f"Document has organization access but user not in org - denying access to {user_email}")
                    return False
            
            # Check direct user permissions
            authorized_users = permission_data.get("users", [])
            if user_email.lower() in [u.lower() for u in authorized_users]:
                print(f"User {user_email} found in direct permissions - allowing access")
                return True
            
            # Check group permissions
            authorized_groups = permission_data.get("groups", [])
            if authorized_groups:
                from app.services.document_permission import check_user_group_membership
                if check_user_group_membership(user_email, authorized_groups, self):
                    print(f"User {user_email} is member of authorized group - allowing access")
                    return True
            
            # No access found
            print(f"User {user_email} not found in authorized users/groups - denying access")
            return False
            
        except Exception as e:
            print(f"Error checking permissions for document {document_id}: {e}")
            # Default to deny access on error
            return False
    def get_document_content(self, document_id: str, user_email: str = None) -> str:
        """
        Get the text content of a document if the user has access.
        
        Args:
            document_id: The document ID
            user_email: The user's email to check permissions
            
        Returns:
            str: The document text content
        """
        token = self.get_access_token()
        drives_endpoint = f"https://graph.microsoft.com/v1.0/sites/{settings.SITE_ID}/drives"
        headers = {"Authorization": f"Bearer {token}"}
        
        # Get drives
        drives_response = requests.get(drives_endpoint, headers=headers)
        drives_response.raise_for_status()
        drives = drives_response.json().get("value", [])
        
        if not drives:
            raise Exception("No drives found in SharePoint response")
        
        drive_id = drives[0]["id"]
        
        # If user_email is provided and not in dev mode, check permissions
        if user_email and not settings.DEV_MODE:
            has_permission = self.check_user_permission(document_id, user_email, drive_id)
            if not has_permission:
                raise Exception(f"User {user_email} does not have permission to access document {document_id}")
        
        # Get document content
        endpoint = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{document_id}/content"
        response = requests.get(endpoint, headers=headers)
        response.raise_for_status()
        
        # Extract text based on content type
        content_type = response.headers.get("Content-Type", "")
        
        if "application/pdf" in content_type:
            with pdfplumber.open(io.BytesIO(response.content)) as pdf:
                text = " ".join(page.extract_text() or "" for page in pdf.pages)
            return text
        elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type:
            text = docx2txt.process(io.BytesIO(response.content))
            return text
        elif "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" in content_type:
            # Excel file
            wb = openpyxl.load_workbook(io.BytesIO(response.content), data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text.append(" ".join(row_text))
            return "\n".join(text)
        elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in content_type:
            # PowerPoint file
            prs = Presentation(io.BytesIO(response.content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif "text/csv" in content_type or response.headers.get("Content-Disposition", "").endswith(".csv"):
            # CSV file
            decoded = response.content.decode("utf-8")
            reader = csv.reader(io.StringIO(decoded))
            text = []
            for row in reader:
                text.append(", ".join(row))
            return "\n".join(text)
        elif "text/plain" in content_type or response.headers.get("Content-Disposition", "").endswith(".txt"):
            # Plain text file
            return response.content.decode("utf-8")
        else:
            return response.text
    
    def get_webhook_subscriptions(self) -> List[Dict[str, Any]]:
        """
        Retrieve all active webhook subscriptions.

        Returns:
            List[Dict[str, Any]]: List of active subscriptions.
        """
        token = self.get_access_token()
        endpoint = "https://graph.microsoft.com/v1.0/subscriptions"
        headers = {"Authorization": f"Bearer {token}"}

        response = requests.get(endpoint, headers=headers)
        response.raise_for_status()
        return response.json().get("value", [])
        
    def create_webhook_subscription(self, resource: str, notification_url: str, expiration_days: int = 3) -> Dict[str, Any]:
        """
        Create a webhook subscription for a SharePoint resource.

        Args:
            resource: The resource to monitor (e.g., drive ID or folder path).
            notification_url: The URL to receive webhook notifications.
            expiration_days: Number of days before the subscription expires (max 3 days).

        Returns:
            Dict[str, Any]: The subscription details.
        """
        token = self.get_access_token()
        endpoint = "https://graph.microsoft.com/v1.0/subscriptions"
        headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
        expiration_date = (datetime.datetime.now(datetime.UTC) + datetime.timedelta(days=3)).replace(microsecond=0).isoformat().replace("+00:00", "Z")

        payload = {
            "changeType": "updated",
            "notificationUrl": notification_url,
            "resource": resource,
            "expirationDateTime": expiration_date,
            "clientState": "secretClientValue"
        }
        
        try:
            response = requests.post(endpoint, headers=headers, json=payload)
            response.raise_for_status()
        except requests.exceptions.RequestException as e:
            print(f"Error creating webhook subscription: {e}")
            if response is not None:
                print(f"Response content: {response.text}")
            raise
    
    def renew_webhook_subscription(self, subscription_id: str, expiration_days: int = 3) -> Dict[str, Any]:
        """
        Renew a webhook subscription for a SharePoint resource.

        Args:
            subscription_id: The ID of the subscription to renew.
            expiration_days: Number of days before the subscription expires (max 3 days).

        Returns:
            Dict[str, Any]: The updated subscription details.
        """
        token = self.get_access_token()
        endpoint = f"https://graph.microsoft.com/v1.0/subscriptions/{subscription_id}"
        headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
        expiration_date = (datetime.datetime.now(datetime.UTC) + datetime.timedelta(days=3)).replace(microsecond=0).isoformat().replace("+00:00", "Z")

        payload = {
            "expirationDateTime": expiration_date
        }

        response = requests.patch(endpoint, headers=headers, json=payload)
        response.raise_for_status()
        return response.json()
        
            
    def process_webhook_notification(self, notification: Dict[str, Any]) -> None:
        """
        Process a webhook notification for document and folder changes.

        Args:
            notification: The webhook notification payload.
        """
        try:
            # Extract the notification value array
            notifications = notification.get("value", [])
            if not notifications:
                print("Invalid notification payload: Missing 'value'")
                return

            for item in notifications:
                resource = item.get("resource")
                change_type = item.get("changeType", "unknown")

                if not resource:
                    print("Invalid notification payload: Missing 'resource'")
                    continue

                print(f"Webhook notification received for resource: {resource}")
                print(f"Change type: {change_type}")

                # Use the Delta API to fetch changes
                token = self.get_access_token()
                delta_link = self.get_delta_link(resource)  # Retrieve the last delta link if available
                endpoint = delta_link or f"https://graph.microsoft.com/v1.0/{resource}/delta"
                headers = {"Authorization": f"Bearer {token}"}
                response = requests.get(endpoint, headers=headers)
                response.raise_for_status()

                delta_response = response.json()
                changes = delta_response.get("value", [])
                if not changes:
                    print("No changes detected in the delta response.")
                    return

                # Save the new delta link for future use
                new_delta_link = delta_response.get("@odata.deltaLink")
                if new_delta_link:
                    self.save_delta_link(resource, new_delta_link)

                # Process each change
                for change in changes:
                    change_type = change.get("@odata.type", "")
                    document_id = change.get("id", "")
                    document_name = change.get("name", "Unknown Name")
                    web_url = change.get("webUrl", "")
                    last_modified = change.get("lastModifiedDateTime", "")
                    drive_id = change.get("parentReference", {}).get("driveId", "")
                    is_folder = "folder" in change  # Check if the item is a folder

                    # Skip the root folder
                    if document_name.lower() == "root":
                        print(f"Skipping root folder: {document_id}")
                        continue

                    print(f"Processing change for ID: {document_id}, Name: {document_name}, Type: {change_type}, Is Folder: {is_folder}")

                    if change_type == "#microsoft.graph.driveItem":
                        if "deleted" in change:
                            print(f"Item deleted: {document_id}, Name: {document_name}")
                            self.delete_document_from_database(document_id)
                        elif is_folder:
                            print(f"Folder created/updated: {document_id}, Name: {document_name}")
                            # Add logic to handle folder changes if needed
                        else:
                            content = self.get_document_content(document_id)
                            if content is None:
                                print(f"Skipping document {document_id} as it could not be fetched.")
                                continue

                            print(f"Document created/updated: {document_id}, Name: {document_name}")
                            self.update_document_in_database(
                                document_id=document_id,
                                document_name=document_name,
                                content=content,
                                drive_id=drive_id,
                                web_url=web_url,
                                last_modified=last_modified,
                            )

        except Exception as e:
            print(f"Error processing webhook notification: {e}")
            
    def get_drive_delta(self, drive_id: str) -> List[Dict[str, Any]]:
        """
        Fetch delta changes for a drive.

        Args:
            drive_id: The ID of the drive.

        Returns:
            List[Dict[str, Any]]: List of changes (created, updated, deleted).
        """
        try:
            token = self.get_access_token()
            endpoint = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/delta"
            headers = {"Authorization": f"Bearer {token}"}

            response = requests.get(endpoint, headers=headers)
            response.raise_for_status()
            return response.json().get("value", [])
        except Exception as e:
            print(f"Error fetching delta changes for drive {drive_id}: {e}")
            return []
        
    def save_delta_link(self, resource: str, delta_link: str) -> None:
        """
        Save the delta link for a resource.

        Args:
            resource: The resource being monitored.
            delta_link: The delta link to save.
        """
        try:
            self.collection.update_one(
                {"resource": resource},
                {"$set": {"delta_link": delta_link}},
                upsert=True
            )
            print(f"Delta link saved for resource: {resource}")
        except Exception as e:
            print(f"Error saving delta link for resource {resource}: {e}")

    def get_delta_link(self, resource: str) -> Optional[str]:
        """
        Retrieve the last saved delta link for a resource.

        Args:
            resource: The resource being monitored.

        Returns:
            Optional[str]: The last saved delta link, or None if not found.
        """
        try:
            record = self.collection.find_one({"resource": resource})
            return record.get("delta_link") if record else None
        except Exception as e:
            print(f"Error retrieving delta link for resource {resource}: {e}")
            return None
    
    def update_document_in_database(self, document_id: str, document_name: str, content: str, drive_id: str, web_url: str, last_modified: str) -> None:
        """
        Update or insert a document in the database, processing it into chunks with metadata and embeddings.

        Args:
            document_id: The ID of the document.
            document_name: The name of the document.
            content: The content of the document.
            drive_id: The ID of the drive containing the document.
            web_url: The URL of the document in SharePoint.
            last_modified: The last modified timestamp of the document.
        """
        try:
            # Initialize text splitter
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

            # Split the document content into chunks
            chunks = text_splitter.split_text(content)

            # Get permissions for the document
            permission_data = get_document_permissions(self, document_id, drive_id)

            # Prepare chunks with metadata
            all_documents = []
            for index, chunk in enumerate(chunks):
                metadata = {
                    "documentId": document_id,
                    "documentName": document_name,
                    "webUrl": web_url,
                    "lastModified": last_modified,
                    "chunkIndex": index,
                    "authorized_users": permission_data.get("users", []),
                    "authorized_groups": permission_data.get("groups", []),
                    "access_level": permission_data.get("access_level", "private"),
                }
                all_documents.append(Document(page_content=chunk, metadata=metadata))

            if not all_documents:
                print(f"No content to update for document {document_name} (ID: {document_id}).")
                return

            # Print sample metadata to verify permissions are included
            print("\nSample document metadata before embedding:")
            for i in range(min(3, len(all_documents))):
                print(f"Document {i+1} metadata: {all_documents[i].metadata}")

            # Get the appropriate embedding model
            embedding_model = get_embedding_model()

            # Remove existing chunks for the document
            self.collection.delete_many({"documentId": document_id})

            # Store documents with embeddings
            MongoDBAtlasVectorSearch.from_documents(
                documents=all_documents,
                embedding=embedding_model,
                collection=self.collection,
                index_name="vector_index",
                text_key="embedding_text",
                embedding_key="embedding",
            )

            print(f"Document {document_id} ({document_name}) updated in the database with {len(all_documents)} chunks.")
        except Exception as e:
            print(f"Error updating document {document_id} in the database: {e}")

    def delete_document_from_database(self, document_id: str) -> None:
        """
        Delete all chunks of a document from the database.

        Args:
            document_id: The ID of the document.
        """
        try:
            result = self.collection.delete_many({"documentId": document_id})
            if result.deleted_count > 0:
                print(f"Document {document_id} deleted from the database.")
            else:
                print(f"Document {document_id} not found in the database.")
        except Exception as e:
            print(f"Error deleting document {document_id} from the database: {e}")
    
   

    def delete_webhook_subscription(self, subscription_id: str):
        """
        Delete a webhook subscription.

        Args:
            subscription_id: The ID of the subscription to delete.
        """
        token = self.get_access_token()
        endpoint = f"https://graph.microsoft.com/v1.0/subscriptions/{subscription_id}"
        headers = {"Authorization": f"Bearer {token}"}
        response = requests.delete(endpoint, headers=headers)
        response.raise_for_status()
        